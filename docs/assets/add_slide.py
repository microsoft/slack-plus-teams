from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

prs = Presentation('C:/source/slack-plus-teams/docs/slack-plus-teams-architecture.pptx')

# Style constants from existing slides
INDIGO = RGBColor(0x4B, 0x4A, 0x78)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT = RGBColor(0x6D, 0x6D, 0x6D)
LIGHT_GRAY = RGBColor(0x7A, 0x7A, 0x7A)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
AMBER = RGBColor(0xEF, 0x6C, 0x00)
RED_ISH = RGBColor(0xC6, 0x28, 0x28)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)

slide_layout = prs.slide_layouts[0]  # only layout available
slide = prs.slides.add_slide(slide_layout)


def add_rect(left, top, width, height, fill_color=None):
    """Add a plain rectangle (squared corners) matching slide 2 header style."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Emu(91440)
    tf.margin_right = Emu(91440)
    tf.margin_top = Emu(45720)
    tf.margin_bottom = Emu(45720)
    return shape, tf


def add_rounded(left, top, width, height, fill_color=None, radius=8000):
    """Add a rounded rectangle with a controlled corner radius."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    # Set corner radius via adjustment
    spPr = shape._element.spPr
    prstGeom = spPr.find(qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            from lxml import etree
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        else:
            for child in list(avLst):
                avLst.remove(child)
        from lxml import etree
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', f'val {radius}')
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Emu(91440)
    tf.margin_right = Emu(91440)
    tf.margin_top = Emu(45720)
    tf.margin_bottom = Emu(45720)
    return shape, tf


def set_text(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.runs[0].font.size = size
    p.runs[0].font.color.rgb = color
    p.runs[0].font.bold = bold


# --- Title banner (squared corners to match slide 2) ---
shape, tf = add_rect(Emu(109728), Emu(411480), Emu(4500000), Emu(713232), fill_color=INDIGO)
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Prioritization Matrix"
run.font.size = Pt(22)
run.font.bold = True
run.font.color.rgb = WHITE

# Subtitle
shape, tf = add_rect(Emu(4800000), Emu(493776), Emu(7000000), Emu(219456))
set_text(tf, "Plot LOB apps by business impact and migration complexity to sequence your rollout.", Pt(13.5), GRAY_TEXT)

# --- Grid layout (smaller, with gaps between quadrants) ---
grid_left = Emu(2200000)
grid_top = Emu(1500000)
grid_w = Emu(7600000)
grid_h = Emu(4200000)
gap = Emu(80000)  # whitespace between quadrants
quad_w = (grid_w - gap) // 2
quad_h = (grid_h - gap) // 2

# Quadrant backgrounds (with subtle corner radius and gap spacing)
add_rounded(grid_left, grid_top, quad_w, quad_h,
            fill_color=RGBColor(0xE8, 0xF5, 0xE9), radius=5000)  # TL green
add_rounded(grid_left + quad_w + gap, grid_top, quad_w, quad_h,
            fill_color=RGBColor(0xFF, 0xF3, 0xE0), radius=5000)  # TR amber
add_rounded(grid_left, grid_top + quad_h + gap, quad_w, quad_h,
            fill_color=LIGHT_BG, radius=5000)  # BL gray
add_rounded(grid_left + quad_w + gap, grid_top + quad_h + gap, quad_w, quad_h,
            fill_color=RGBColor(0xFF, 0xEB, 0xEE), radius=5000)  # BR red


# Quadrant labels
def add_quadrant_label(left, top, text, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + Emu(80000), top + Emu(60000),
        Emu(1800000), Emu(280000)
    )
    shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    p.runs[0].font.size = Pt(11)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = color


add_quadrant_label(grid_left, grid_top, "Migrate First", GREEN)
add_quadrant_label(grid_left + quad_w + gap, grid_top, "Plan Carefully", AMBER)
add_quadrant_label(grid_left, grid_top + quad_h + gap, "Quick Wins", RGBColor(0x75, 0x75, 0x75))
add_quadrant_label(grid_left + quad_w + gap, grid_top + quad_h + gap, "Deprioritize", RED_ISH)


# --- App bubbles (smaller, reduced corner radius) ---
def add_app_bubble(cx, cy, name, subtitle, fill_color, text_color=WHITE):
    bw = Emu(1700000)
    bh = Emu(560000)
    shape, tf = add_rounded(cx - bw // 2, cy - bh // 2, bw, bh,
                            fill_color=fill_color, radius=5000)
    tf.margin_left = Emu(64000)
    tf.margin_right = Emu(64000)
    tf.margin_top = Emu(32000)
    tf.margin_bottom = Emu(32000)
    p = tf.paragraphs[0]
    p.text = name
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(11)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = text_color
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.alignment = PP_ALIGN.CENTER
    p2.runs[0].font.size = Pt(9)
    p2.runs[0].font.color.rgb = text_color
    p2.runs[0].font.bold = False


# Quadrant centers (based on new smaller grid)
tl_cx = grid_left + quad_w // 2
tl_cy = grid_top + quad_h // 2
tr_cx = grid_left + quad_w + gap + quad_w // 2
tr_cy = grid_top + quad_h // 2
bl_cx = grid_left + quad_w // 2
bl_cy = grid_top + quad_h + gap + quad_h // 2
br_cx = grid_left + quad_w + gap + quad_w // 2
br_cy = grid_top + quad_h + gap + quad_h // 2

# Top-left: High Impact, Low Complexity
add_app_bubble(tl_cx, tl_cy - Emu(280000), "IT Helpdesk Bot", "Q&A, ticket creation", GREEN)
add_app_bubble(tl_cx, tl_cy + Emu(280000), "Sales Dashboard Bot", "CRM lookups, daily reports", GREEN)

# Top-right: High Impact, High Complexity
add_app_bubble(tr_cx, tr_cy, "HR Onboarding Suite", "Workflows, forms, approvals", AMBER)

# Bottom-left: Low Impact, Low Complexity
add_app_bubble(bl_cx, bl_cy, "Standup Reminder", "Scheduled messages only", RGBColor(0x75, 0x75, 0x75))

# Bottom-right: Low Impact, High Complexity
add_app_bubble(br_cx, br_cy, "Legacy Build Monitor", "Custom integrations, webhooks", RED_ISH)


# --- Axis labels ---
def add_label(left, top, width, height, text, size, color, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.runs[0].font.size = size
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = color


# Y-axis labels
add_label(Emu(400000), grid_top + Emu(500000), Emu(1400000), Emu(350000),
          "High Impact", Pt(10), INDIGO)
add_label(Emu(400000), grid_top + quad_h + gap + Emu(500000), Emu(1400000), Emu(350000),
          "Low Impact", Pt(10), INDIGO)

# X-axis labels
total_w = quad_w * 2 + gap
add_label(grid_left + Emu(400000), grid_top + grid_h + Emu(80000),
          Emu(2000000), Emu(300000), "Low Complexity", Pt(10), INDIGO)
add_label(grid_left + quad_w + gap + Emu(400000), grid_top + grid_h + Emu(80000),
          Emu(2000000), Emu(300000), "High Complexity", Pt(10), INDIGO)

# Axis titles
add_label(Emu(400000), grid_top + quad_h // 2, Emu(1200000), Emu(350000),
          "Business Impact \u2191", Pt(9), INDIGO)
add_label(grid_left + quad_w - Emu(200000), grid_top + grid_h + Emu(340000),
          Emu(1800000), Emu(250000), "Migration Complexity \u2192", Pt(9), INDIGO)

# --- Footer ---
shape, tf = add_rect(Emu(219456), Emu(6455664), Emu(6035040), Emu(128016))
set_text(tf, "Slack to Teams migration guide \u2022 architecture diagram 3 of 7",
         Pt(9.5), LIGHT_GRAY)

prs.save('C:/source/slack-plus-teams/docs/slack-plus-teams-architecture-v2.pptx')
print("Slide 3 (Prioritization Matrix) added successfully!")
