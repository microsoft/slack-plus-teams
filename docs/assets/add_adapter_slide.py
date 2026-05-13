from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

prs = Presentation('C:/source/slack-plus-teams/docs/slack-plus-teams-architecture.pptx')

# Theme colors from slides 1 and 2 — no platform-specific colors
INDIGO = RGBColor(0x4B, 0x4A, 0x78)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT = RGBColor(0x6D, 0x6D, 0x6D)
LIGHT_GRAY = RGBColor(0x7A, 0x7A, 0x7A)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
LIGHT_INDIGO_BG = RGBColor(0xE8, 0xE8, 0xF0)  # light indigo for adapter bg
LIGHT_GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)    # light green for shared layer bg
ARROW_COLOR = INDIGO

slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)


def add_rect(left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Emu(91440)
    tf.margin_right = Emu(91440)
    tf.margin_top = Emu(45720)
    tf.margin_bottom = Emu(45720)
    return shape, tf


def add_rounded(left, top, width, height, fill_color=None, radius=5000,
                border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(1.5)
    else:
        shape.line.fill.background()
    spPr = shape._element.spPr
    prstGeom = spPr.find(qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            from lxml import etree
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        else:
            for child in list(avLst):
                avLst.remove(child)
        from lxml import etree
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', f'val {radius}')
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Emu(91440)
    tf.margin_right = Emu(91440)
    tf.margin_top = Emu(45720)
    tf.margin_bottom = Emu(45720)
    return shape, tf


def set_text(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.runs[0].font.size = size
    p.runs[0].font.color.rgb = color
    p.runs[0].font.bold = bold


def add_box(left, top, width, height, title, subtitle_lines,
            fill_color, text_color, border_color=None):
    """Add a labeled box with title and subtitle lines."""
    shape, tf = add_rounded(left, top, width, height, fill_color=fill_color,
                            radius=4000, border_color=border_color,
                            border_width=Pt(1.5))
    tf.margin_top = Emu(36000)
    tf.margin_bottom = Emu(36000)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = text_color
    for line in subtitle_lines:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = line
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = text_color
        r2.font.bold = False
    return shape


def add_arrow(x1, y1, x2, y2, color=ARROW_COLOR, width=Pt(2), head=True, tail=False):
    """Add an arrow. head=arrowhead at end, tail=arrowhead at start."""
    connector = slide.shapes.add_connector(
        1, Emu(x1), Emu(y1), Emu(x2), Emu(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = width
    connector.line.fill.solid()
    connector.line.color.rgb = color
    ln = connector._element.find(qn('a:ln'))
    if ln is None:
        spPr = connector._element.spPr
        ln = spPr.find(qn('a:ln'))
    if ln is not None:
        from lxml import etree
        if head:
            t = etree.SubElement(ln, qn('a:tailEnd'))
            t.set('type', 'triangle')
            t.set('w', 'med')
            t.set('len', 'med')
        if tail:
            h = etree.SubElement(ln, qn('a:headEnd'))
            h.set('type', 'triangle')
            h.set('w', 'med')
            h.set('len', 'med')
    return connector


# --- Title banner (squared corners, matching slide 2) ---
shape, tf = add_rect(Emu(109728), Emu(411480), Emu(4500000), Emu(713232), fill_color=INDIGO)
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Adapter Pattern Architecture"
run.font.size = Pt(22)
run.font.bold = True
run.font.color.rgb = WHITE

# Subtitle
shape, tf = add_rect(Emu(4800000), Emu(493776), Emu(7200000), Emu(219456))
set_text(tf, "Inbound calls from each platform converge on a shared service layer",
         Pt(13.5), GRAY_TEXT)


# === Layout ===
#
#   [Slack]  -->  [Slack Adapter]  -->  [  Shared Service  ]  <--  [Teams Adapter]  <--  [Teams]
#                                       [      Layer       ]
#
# Both platforms send calls INWARD to the central shared layer.
# The shared layer is the hub; adapters translate platform-specific protocols.

slide_w = 12192000
center_x = slide_w // 2
center_y = 3700000  # slightly below center to account for title

# Shared Service Layer — center, prominent
shared_w = 2600000
shared_h = 2000000
shared_left = center_x - shared_w // 2
shared_top = center_y - shared_h // 2

# Adapter boxes
adapter_w = 1800000
adapter_h = 1400000
adapter_gap = 400000  # gap between adapter and shared layer

slack_adapter_left = shared_left - adapter_gap - adapter_w
slack_adapter_top = center_y - adapter_h // 2

teams_adapter_left = shared_left + shared_w + adapter_gap
teams_adapter_top = center_y - adapter_h // 2

# Platform boxes (outer)
platform_w = 1500000
platform_h = 1000000
platform_gap = 350000

slack_left = slack_adapter_left - platform_gap - platform_w
slack_top = center_y - platform_h // 2

teams_left = teams_adapter_left + adapter_w + platform_gap
teams_top = center_y - platform_h // 2


# --- Draw boxes ---

# Shared Service Layer (center, green theme)
add_box(shared_left, shared_top, shared_w, shared_h,
        "Shared Service Layer",
        ["Business logic", "AI / LLM calls", "Data access", "Common types"],
        fill_color=LIGHT_GREEN_BG, text_color=GREEN,
        border_color=GREEN)

# Slack Adapter (left of center)
add_box(slack_adapter_left, slack_adapter_top, adapter_w, adapter_h,
        "Slack Adapter",
        ["@slack/bolt", "Event handlers", "Block Kit"],
        fill_color=LIGHT_INDIGO_BG, text_color=INDIGO,
        border_color=INDIGO)

# Teams Adapter (right of center)
add_box(teams_adapter_left, teams_adapter_top, adapter_w, adapter_h,
        "Teams Adapter",
        ["Agents SDK", "Activity handlers", "Adaptive Cards"],
        fill_color=LIGHT_INDIGO_BG, text_color=INDIGO,
        border_color=INDIGO)

# Slack Platform (far left)
add_box(slack_left, slack_top, platform_w, platform_h,
        "Slack",
        ["WebSocket", "Socket Mode"],
        fill_color=INDIGO, text_color=WHITE)

# Teams Platform (far right)
add_box(teams_left, teams_top, platform_w, platform_h,
        "Microsoft Teams",
        ["HTTPS", "Bot Framework"],
        fill_color=INDIGO, text_color=WHITE)


# --- Arrows: both sides point INWARD to the shared layer ---
arrow_y_top = center_y - 150000    # upper arrow (inbound)
arrow_y_bot = center_y + 150000    # lower arrow (outbound/response)

# LEFT SIDE: Slack → Slack Adapter → Shared Layer

# Slack → Slack Adapter (inbound)
add_arrow(slack_left + platform_w, arrow_y_top,
          slack_adapter_left, arrow_y_top)

# Slack Adapter → Shared Layer (inbound)
add_arrow(slack_adapter_left + adapter_w, arrow_y_top,
          shared_left, arrow_y_top)

# Shared Layer → Slack Adapter (response)
add_arrow(shared_left, arrow_y_bot,
          slack_adapter_left + adapter_w, arrow_y_bot)

# Slack Adapter → Slack (response)
add_arrow(slack_adapter_left, arrow_y_bot,
          slack_left + platform_w, arrow_y_bot)


# RIGHT SIDE: Teams → Teams Adapter → Shared Layer

# Teams → Teams Adapter (inbound)
add_arrow(teams_left, arrow_y_top,
          teams_adapter_left + adapter_w, arrow_y_top)

# Teams Adapter → Shared Layer (inbound)
add_arrow(teams_adapter_left, arrow_y_top,
          shared_left + shared_w, arrow_y_top)

# Shared Layer → Teams Adapter (response)
add_arrow(shared_left + shared_w, arrow_y_bot,
          teams_adapter_left, arrow_y_bot)

# Teams Adapter → Teams (response)
add_arrow(teams_adapter_left + adapter_w, arrow_y_bot,
          teams_left, arrow_y_bot)


# --- Arrow labels ---
def add_label(cx, cy, text):
    w = 1400000
    h = 220000
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Emu(cx - w // 2), Emu(cy),
                                   Emu(w), Emu(h))
    shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.runs[0].font.size = Pt(8.5)
    p.runs[0].font.color.rgb = LIGHT_GRAY
    p.runs[0].font.italic = True


# Labels above the inbound arrows
label_y = center_y - shared_h // 2 - 280000

# Left side labels
slack_mid = (slack_left + platform_w + slack_adapter_left) // 2
add_label(slack_mid, label_y, "Events / Messages")

slack_adapter_mid = (slack_adapter_left + adapter_w + shared_left) // 2
add_label(slack_adapter_mid, label_y, "Normalized Calls")

# Right side labels
teams_mid = (teams_left + teams_adapter_left + adapter_w) // 2
add_label(teams_mid, label_y, "Activities / Messages")

teams_adapter_mid = (teams_adapter_left + shared_left + shared_w) // 2
add_label(teams_adapter_mid, label_y, "Normalized Calls")


# --- Footer ---
shape, tf = add_rect(Emu(219456), Emu(6455664), Emu(6035040), Emu(128016))
set_text(tf, "Slack to Teams migration guide \u2022 architecture diagram 4 of 7",
         Pt(9.5), LIGHT_GRAY)

prs.save('C:/source/slack-plus-teams/docs/slack-plus-teams-architecture-v3.pptx')
print("Slide 4 (Adapter Pattern Architecture) added successfully!")
